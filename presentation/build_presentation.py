"""
Generate a PowerPoint (.pptx) presentation for the AI Mutual Fund Analyzer project.

Requirements (install before running):
- python-pptx
- pillow

Optional (for diagram rendering):
- cairosvg (convert SVG to PNG)
- mermaid-cli (render mermaid to SVG) — alternatively render diagrams using an online renderer

Usage:
python presentation/build_presentation.py --output presentation/AI_Mutual_Fund_Presentation.pptx

If you want a PDF, convert the generated PPTX using LibreOffice:
soffice --headless --convert-to pdf presentation/AI_Mutual_Fund_Presentation.pptx --outdir presentation/

"""
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception as e:
    print("Missing dependency: python-pptx. Install with `pip install python-pptx pillow` and try again.")
    raise

BASE = Path(__file__).resolve().parent
SLIDE_MD = BASE / "slide_content.md"
DIAGRAMS_DIR = BASE / "diagrams"


def read_slides(md_path):
    slides = []
    if not md_path.exists():
        return slides
    text = md_path.read_text(encoding="utf-8")
    sections = text.split('\n\n---\n\n')
    for sec in sections:
        lines = [l.rstrip() for l in sec.strip().splitlines() if l.strip()]
        if not lines:
            continue
        title = lines[0].replace('Title: ', '').replace('Slide: ', '')
        content = '\n'.join(lines[1:])
        slides.append((title, content))
    return slides


def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        try:
            slide.placeholders[1].text = subtitle
        except Exception:
            pass


def add_content_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    for line in content.splitlines():
        if line.strip().startswith('- '):
            p = body.add_paragraph() if body.paragraphs else body.paragraphs[0]
            p.level = 0
            p.text = line.strip()[2:]
            p.font.size = Pt(18)
        elif line.strip().startswith('1.') or line.strip().startswith('2.'):
            p = body.add_paragraph()
            p.text = line.strip()
            p.level = 0
            p.font.size = Pt(18)
        else:
            p = body.add_paragraph()
            p.text = line.strip()
            p.level = 0
            p.font.size = Pt(18)


def add_diagram_slide(prs, title, mermaid_file):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.7)
    width = Inches(9)
    # Try to find a PNG with same name
    png_path = mermaid_file.with_suffix('.png')
    if png_path.exists():
        slide.shapes.add_picture(str(png_path), left, top, width=width)
        return
    # If no PNG, insert the mermaid source as text
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    body.text = 'Diagram (render separately if you want an image):\n\n'
    body.add_paragraph().text = mermaid_file.read_text(encoding='utf-8')


def build_presentation(output_path):
    prs = Presentation()
    slides = read_slides(SLIDE_MD)
    if not slides:
        add_title_slide(prs, 'AI Mutual Fund Analyzer & Recommender System', 'End-to-end project overview')
    for i, (title, content) in enumerate(slides):
        if i == 0 and title.lower().startswith('title'):
            # parse title slide
            # we expect first lines like Title: ... and Subtitle: ...
            lines = content.splitlines()
            subtitle = None
            if lines and lines[0].lower().startswith('subtitle:'):
                subtitle = lines[0].split(':', 1)[1].strip()
            add_title_slide(prs, title.replace('Title: ', ''), subtitle)
            continue
        if 'Architecture Diagram' in title or 'Architecture' in title:
            mermaid_file = DIAGRAMS_DIR / 'architecture.mmd'
            add_diagram_slide(prs, title, mermaid_file)
        else:
            add_content_slide(prs, title, content)
    prs.save(output_path)
    print(f"Saved presentation to {output_path}")


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument('--output', '-o', default=str(BASE / 'AI_Mutual_Fund_Presentation.pptx'))
    args = parser.parse_args()
    out = Path(args.output)
    out.parent.mkdir(parents=True, exist_ok=True)
    build_presentation(out)
    print('Done. Convert to PDF with LibreOffice if needed.')
